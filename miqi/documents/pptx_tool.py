"""PowerPoint (.pptx) read/write tools."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from miqi.agent.tools.base import Tool


def _raw_output_path(kwargs: dict[str, Any]) -> str:
    return str(
        kwargs.get("filename")
        or kwargs.get("file_path")
        or kwargs.get("path")
        or ""
    )


def _ensure_suffix(path: Path, suffix: str) -> Path:
    if not path.name or path.name in {".", ".."}:
        raise ValueError("output filename is required")
    if path.suffix.lower() == suffix:
        return path
    return path.with_suffix(suffix)


def _enforce_boundary(path: Path, allowed_dir: Path | None, workspace: Path | None) -> None:
    effective_dir = allowed_dir or workspace
    if effective_dir is None:
        return
    try:
        path.resolve().relative_to(effective_dir.resolve())
    except ValueError:
        raise PermissionError(
            f"Path '{path}' resolves outside allowed directory '{effective_dir}'"
        )


def _resolve_output_path(
    file_path: str,
    workspace: Path | None,
    allowed_dir: Path | None,
) -> Path:
    """Resolve an output path and enforce workspace/directory bounds.

    Office document write tools always write inside the workspace:
    - Relative paths are resolved against *workspace*.
    - If *allowed_dir* is ``None`` but *workspace* is set, *workspace*
      is used as the effective boundary (defense-in-depth default).
    - Absolute paths outside the effective boundary are rejected.

    Raises:
        PermissionError: if the resolved path is outside the effective boundary.
    """
    p = Path(file_path).expanduser()
    if not p.is_absolute() and workspace is not None:
        p = workspace / p
    resolved = p.resolve()

    # Defense-in-depth: when no explicit allowed_dir is given, office
    # write tools default to workspace as the boundary.
    effective_dir = allowed_dir
    if effective_dir is None and workspace is not None:
        effective_dir = workspace.resolve()

    if effective_dir is not None:
        try:
            resolved.relative_to(effective_dir.resolve())
        except ValueError:
            raise PermissionError(
                f"Path '{file_path}' resolves outside allowed directory "
                f"'{effective_dir}'"
            )
    return resolved


class PptxReadTool(Tool):
    """Read the content of a PowerPoint (.pptx) presentation."""

    name = "pptx_read"
    description = "Read and extract text content from a PowerPoint (.pptx) file."

    def __init__(
        self,
        workspace: Path | None = None,
        allowed_dir: Path | None = None,
    ):
        self._workspace = workspace
        self._allowed_dir = allowed_dir

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path to the .pptx file to read. Alias for filename.",
                },
                "filename": {
                    "type": "string",
                    "description": "Filename or relative path to the .pptx file.",
                },
                "path": {
                    "type": "string",
                    "description": "Path to the .pptx file to read. Alias for filename.",
                },
            },
            "anyOf": [
                {"required": ["filename"]},
                {"required": ["file_path"]},
                {"required": ["path"]},
            ],
        }

    async def execute(self, **kwargs: Any) -> str:
        raw_path = _raw_output_path(kwargs)
        if not raw_path.strip():
            return "Error: filename is required"
        try:
            file_path = _resolve_output_path(
                raw_path, self._workspace, self._allowed_dir,
            )
            file_path = _ensure_suffix(file_path, ".pptx")
            _enforce_boundary(file_path, self._allowed_dir, self._workspace)
        except PermissionError as e:
            return f"Error: Permission denied: {e}"
        except ValueError as e:
            return f"Error: {e}"
        if not file_path.exists():
            return f"Error: file not found: {file_path}"
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            lines = [f"Presentation: {len(prs.slides)} slides", ""]
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)
                lines.append("")
            return "\n".join(lines)
        except Exception as e:
            return f"Error reading {file_path.name}: {e}"


class CreatePptxTool(Tool):
    """Create a PowerPoint (.pptx) presentation."""

    name = "create_pptx"
    description = (
        "Create a PowerPoint (.pptx) presentation in the workspace files directory. "
        "Supports multiple slides with titles, bullets, body text, and images."
    )

    def __init__(
        self,
        workspace: Path | None = None,
        allowed_dir: Path | None = None,
    ):
        self._workspace = workspace
        self._allowed_dir = allowed_dir

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "Path for the output .pptx file. Alias for filename.",
                },
                "filename": {
                    "type": "string",
                    "description": "Filename or relative path for the output .pptx file",
                },
                "path": {
                    "type": "string",
                    "description": "Path for the output .pptx file. Alias for filename.",
                },
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "subtitle": {"type": "string"},
                            "content": {
                                "description": "Slide body text or an array of body lines.",
                            },
                            "bullets": {"type": "array", "items": {"type": "string"}},
                            "image_path": {"type": "string"},
                        },
                    },
                    "description": "List of slides with title, content, bullets, or image_path",
                },
            },
            "anyOf": [
                {"required": ["filename"]},
                {"required": ["file_path"]},
                {"required": ["path"]},
            ],
        }

    async def execute(self, **kwargs: Any) -> str:
        from pptx import Presentation
        from pptx.util import Inches

        raw_path = _raw_output_path(kwargs)
        slides = kwargs.get("slides") or []
        if not raw_path.strip():
            return "Error: filename is required"

        try:
            file_path = _resolve_output_path(
                raw_path, self._workspace, self._allowed_dir,
            )
            file_path = _ensure_suffix(file_path, ".pptx")
            _enforce_boundary(file_path, self._allowed_dir, self._workspace)
        except PermissionError as e:
            return f"Error: Permission denied: {e}"
        except ValueError as e:
            return f"Error: {e}"
        if not slides:
            return "Error: slides is required"

        try:
            prs = Presentation()
            for slide_data in slides:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                if title:
                    title.text = str(slide_data.get("title", ""))
                body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
                if body and hasattr(body, "text_frame"):
                    text_frame = body.text_frame
                    text_frame.clear()
                    content_items: list[Any] = []
                    if slide_data.get("subtitle"):
                        content_items.append(slide_data["subtitle"])
                    content = slide_data.get("content")
                    if isinstance(content, list):
                        content_items.extend(content)
                    elif content:
                        content_items.append(content)
                    bullets = slide_data.get("bullets") or []
                    for item_index, item in enumerate(content_items):
                        if item_index == 0:
                            text_frame.text = str(item)
                        else:
                            paragraph = text_frame.add_paragraph()
                            paragraph.text = str(item)
                            paragraph.level = 0
                    for bullet in bullets:
                        paragraph = text_frame.add_paragraph()
                        paragraph.text = str(bullet)
                        paragraph.level = 0
                image_path = slide_data.get("image_path")
                if image_path:
                    slide.shapes.add_picture(
                        str(image_path),
                        Inches(float(slide_data.get("image_left", 5.5))),
                        Inches(float(slide_data.get("image_top", 1.5))),
                        width=Inches(float(slide_data.get("image_width", 3.0))),
                    )

            file_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(file_path))
            return f"Created: {file_path} ({len(slides)} slides)"
        except Exception as e:
            return f"Error writing {raw_path}: {e}"


class PptxWriteTool(CreatePptxTool):
    """Backward-compatible alias for create_pptx."""

    name = "pptx_write"
    description = "Create a new PowerPoint (.pptx) file. Prefer create_pptx for new calls."
