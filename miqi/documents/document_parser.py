"""Document parsing service for PDF, Word, PowerPoint, Excel, Markdown, and HTML files.

Provides text extraction for preview and LLM context. PDF parsing supports
OCR fallback for scanned/image-based documents via Tesseract.
Chart and table extraction via pdfplumber for structured data.
HTML parsing via lxml with stdlib fallback.
"""

from __future__ import annotations

import base64
import io
import os
import re
import shutil
import subprocess
import tempfile
import time
from pathlib import Path
from typing import Any

from loguru import logger

try:
    from lxml import html as _lxml_html
    _HAS_LXML_HTML = True
except ImportError:
    _HAS_LXML_HTML = False


# ── Configuration ──────────────────────────────────────────────────────────

# Maximum extracted text length returned to frontend preview
MAX_PREVIEW_CHARS = 50_000

# Maximum extracted text length for LLM context
MAX_CONTEXT_CHARS = 200_000


# ── Public API ─────────────────────────────────────────────────────────────

def parse_document(
    file_path: Path,
    *,
    max_chars: int = MAX_CONTEXT_CHARS,
    force_ocr: bool = False,
    extract_charts: bool = True,
) -> dict[str, Any]:
    """Parse a document file and return extracted text with metadata.

    Args:
        file_path: Path to the document file.
        max_chars: Maximum characters to return.
        force_ocr: If True, force OCR even for text-based PDFs.
        extract_charts: If True, also extract structured tables/charts from PDFs/PPTX.

    Returns:
        dict with keys: text, page_count, size_bytes, mime_type, ocr_used,
                        parse_ms, charts (if extract_charts=True)
    """
    suffix = _get_suffix(file_path)
    if suffix in _PDF_SUFFIXES:
        return _parse_pdf(file_path, max_chars=max_chars, force_ocr=force_ocr,
                          extract_charts=extract_charts)
    elif suffix in _IMAGE_SUFFIXES:
        return _parse_image(file_path, max_chars=max_chars)
    elif suffix in _DOCX_SUFFIXES:
        return _parse_docx(file_path, max_chars=max_chars)
    elif suffix in _PPTX_SUFFIXES:
        return _parse_pptx(file_path, max_chars=max_chars, extract_charts=extract_charts)
    elif suffix in _XLSX_SUFFIXES:
        return _parse_xlsx(file_path, max_chars=max_chars)
    elif suffix in _MD_SUFFIXES:
        return _parse_markdown(file_path, max_chars=max_chars)
    elif suffix in _HTML_SUFFIXES:
        return _parse_html(file_path, max_chars=max_chars)
    elif suffix in _CSV_SUFFIXES:
        return _parse_csv(file_path, max_chars=max_chars)
    elif suffix in _JSON_SUFFIXES:
        return _parse_json(file_path, max_chars=max_chars)
    elif suffix in _XML_FILE_SUFFIXES:
        return _parse_xml_file(file_path, max_chars=max_chars)
    elif suffix in _YAML_SUFFIXES:
        return _parse_yaml(file_path, max_chars=max_chars)
    elif suffix in _ENV_SUFFIXES:
        return _parse_env(file_path, max_chars=max_chars)
    elif suffix in _LOG_SUFFIXES:
        return _parse_log(file_path, max_chars=max_chars)
    elif suffix in _SQL_SUFFIXES:
        return _parse_sql(file_path, max_chars=max_chars)
    elif suffix in _INI_SUFFIXES:
        return _parse_ini(file_path, max_chars=max_chars)
    elif suffix in _TOML_SUFFIXES:
        return _parse_toml(file_path, max_chars=max_chars)
    elif suffix in _HTACCESS_SUFFIXES:
        return _parse_htaccess(file_path, max_chars=max_chars)
    elif suffix in _SH_SUFFIXES:
        return _parse_sh(file_path, max_chars=max_chars)
    elif suffix in _TXT_SUFFIXES:
        return _parse_txt(file_path, max_chars=max_chars)
    elif suffix in _RTF_SUFFIXES:
        return _parse_rtf(file_path, max_chars=max_chars)
    else:
        raise ValueError(f"Unsupported document format: {suffix}")


def is_supported_document(path: Path | str) -> bool:
    """Check if the file path is a supported document format."""
    suffix = _get_suffix(path)
    return suffix in _ALL_DOCUMENT_SUFFIXES


def get_document_category(path: Path | str) -> str:
    """Get the document category (pdf, word, ppt, excel, markdown, unknown)."""
    suffix = _get_suffix(path)
    if suffix in _PDF_SUFFIXES:
        return "pdf"
    if suffix in _DOCX_SUFFIXES:
        return "word"
    if suffix in _PPTX_SUFFIXES:
        return "ppt"
    if suffix in _XLSX_SUFFIXES:
        return "excel"
    if suffix in _MD_SUFFIXES:
        return "markdown"
    if suffix in _HTML_SUFFIXES:
        return "html"
    if suffix in _CSV_SUFFIXES:
        return "csv"
    if suffix in _JSON_SUFFIXES:
        return "json"
    if suffix in _XML_FILE_SUFFIXES:
        return "xml"
    if suffix in _YAML_SUFFIXES:
        return "yaml"
    if suffix in _ENV_SUFFIXES:
        return "env"
    if suffix in _LOG_SUFFIXES:
        return "log"
    if suffix in _SQL_SUFFIXES:
        return "sql"
    if suffix in _INI_SUFFIXES:
        return "ini"
    if suffix in _TOML_SUFFIXES:
        return "toml"
    if suffix in _HTACCESS_SUFFIXES:
        return "htaccess"
    if suffix in _SH_SUFFIXES:
        return "sh"
    if suffix in _TXT_SUFFIXES:
        return "txt"
    if suffix in _RTF_SUFFIXES:
        return "rtf"
    return "unknown"


# ── Suffix maps ────────────────────────────────────────────────────────────

_PDF_SUFFIXES = {".pdf"}
_DOCX_SUFFIXES = {".docx", ".doc", ".odt"}
_PPTX_SUFFIXES = {".pptx", ".ppt", ".odp"}
_XLSX_SUFFIXES = {".xlsx", ".xls", ".ods"}
_MD_SUFFIXES = {".md", ".markdown", ".mdown"}
_HTML_SUFFIXES = {".html", ".htm"}
_CSV_SUFFIXES = {".csv"}
_JSON_SUFFIXES = {".json"}
_XML_FILE_SUFFIXES = {".xml"}
_YAML_SUFFIXES = {".yaml", ".yml"}
_ENV_SUFFIXES = {".env"}
_LOG_SUFFIXES = {".log"}
_SQL_SUFFIXES = {".sql"}
_INI_SUFFIXES = {".ini"}
_TOML_SUFFIXES = {".toml"}
_HTACCESS_SUFFIXES = {".htaccess"}
_SH_SUFFIXES = {".sh", ".bash"}
_TXT_SUFFIXES = {".txt", ".text"}
_RTF_SUFFIXES = {".rtf"}

_IMAGE_SUFFIXES = {
    ".jpg", ".jpeg", ".png", ".bmp", ".gif", ".webp", ".tiff", ".tif", ".ico",
    # SVG 矢量图 — graph_render 产物（#715）：需进 _ALL_DOCUMENT_SUFFIXES
    # 否则后缀校验在查 MIME 前就拒绝（CodeRabbit #761）
    ".svg",
}

_ALL_DOCUMENT_SUFFIXES = (
    _PDF_SUFFIXES | _DOCX_SUFFIXES | _PPTX_SUFFIXES |
    _XLSX_SUFFIXES | _MD_SUFFIXES | _HTML_SUFFIXES |
    _CSV_SUFFIXES | _JSON_SUFFIXES | _XML_FILE_SUFFIXES |
    _YAML_SUFFIXES | _ENV_SUFFIXES | _LOG_SUFFIXES |
    _SQL_SUFFIXES | _INI_SUFFIXES | _TOML_SUFFIXES |
    _HTACCESS_SUFFIXES | _SH_SUFFIXES | _TXT_SUFFIXES |
    _RTF_SUFFIXES | _IMAGE_SUFFIXES
)


# ── Suffix normalization ──────────────────────────────────────────────────

def _get_suffix(file_path: Path | str) -> str:
    """Normalize suffix for dotfiles where Path.suffix returns empty.

    Path(".env").suffix == ""  → _get_suffix(".env") == ".env"
    Path(".htaccess").suffix == ""  → _get_suffix(".htaccess") == ".htaccess"
    Path("test.csv").suffix == ".csv"  → unchanged.
    """
    p = Path(file_path)
    suffix = p.suffix.lower()
    if suffix:
        return suffix
    # Dotfile: name IS the suffix (e.g., ".env", ".htaccess")
    name = p.name.lower()
    if name.startswith("."):
        return name
    return ""


# ── MIME types ─────────────────────────────────────────────────────────────

_SUFFIX_TO_MIME: dict[str, str] = {
    ".pdf": "application/pdf",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".bmp": "image/bmp",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".tiff": "image/tiff",
    ".tif": "image/tiff",
    ".ico": "image/x-icon",
    ".svg": "image/svg+xml",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".doc": "application/msword",
    ".ppt": "application/vnd.ms-powerpoint",
    ".xls": "application/vnd.ms-excel",
    ".odt": "application/vnd.oasis.opendocument.text",
    ".odp": "application/vnd.oasis.opendocument.presentation",
    ".ods": "application/vnd.oasis.opendocument.spreadsheet",
    ".html": "text/html",
    ".htm": "text/html",
    ".csv": "text/csv",
    ".json": "application/json",
    ".xml": "application/xml",
    ".yaml": "text/yaml",
    ".yml": "text/yaml",
    ".env": "text/plain",
    ".log": "text/plain",
    ".sql": "text/x-sql",
    ".ini": "text/plain",
    ".toml": "text/plain",
    ".htaccess": "text/plain",
    ".sh": "text/x-shellscript",
    ".bash": "text/x-shellscript",
    ".txt": "text/plain",
    ".text": "text/plain",
    ".rtf": "text/rtf",
}


# ── PDF Parsing ────────────────────────────────────────────────────────────

def _parse_pdf(
    file_path: Path,
    *,
    max_chars: int = MAX_CONTEXT_CHARS,
    force_ocr: bool = False,
    extract_charts: bool = True,
) -> dict[str, Any]:
    """Extract text from a PDF file with OCR fallback.

    Uses pypdfium2 for fast text extraction first, falls back to pypdf
    if not available, and finally to OCR via Tesseract for scanned PDFs.
    """
    import time
    t0 = time.monotonic()

    text = ""
    page_count = 0
    ocr_used = False

    # Method 1: Try pypdfium2 (fastest — ~1ms per page)
    if not force_ocr:
        try:
            import pypdfium2 as pdfium
            pdf_doc = pdfium.PdfDocument(str(file_path))
            page_count = len(pdf_doc)
            pages_text = []
            for i in range(page_count):
                try:
                    page = pdf_doc[i]
                    textpage = page.get_textpage()
                    page_text = textpage.get_text_range()
                    if page_text and page_text.strip():
                        pages_text.append(page_text)
                except Exception:
                    try:
                        textpage = page.get_textpage()
                        page_text = textpage.get_text_bounded()
                        if page_text and page_text.strip():
                            pages_text.append(page_text)
                    except Exception:
                        pass
            text = "\n\n".join(pages_text)
            pdf_doc.close()
        except ImportError:
            pass
        except Exception as exc:
            logger.warning(f"pypdfium2 extraction failed: {exc}")

    # Method 2: Try pypdf (slower, fallback)
    if not force_ocr and len(text.strip()) < 100:
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            if page_count == 0:
                page_count = len(reader.pages)
            pages_text = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    pages_text.append(page_text)
            text = "\n\n".join(pages_text)
        except ImportError:
            logger.warning("pypdf not installed, falling back to OCR")
        except Exception as exc:
            logger.warning(f"pypdf extraction failed: {exc}")

    # Method 3: OCR via Tesseract if text extraction is insufficient (scanned PDF)
    if force_ocr or len(text.strip()) < 100:
        logger.info(f"PDF text too short ({len(text)} chars, page_count={page_count}), attempting OCR")
        ocr_text = _pdf_ocr(file_path)
        if ocr_text and len(ocr_text) > len(text):
            text = ocr_text
            ocr_used = True

    text = text[:max_chars]
    parse_ms = (time.monotonic() - t0) * 1000

    result: dict[str, Any] = {
        "text": text,
        "page_count": page_count,
        "size_bytes": file_path.stat().st_size,
        "mime_type": "application/pdf",
        "ocr_used": ocr_used,
        "parse_ms": round(parse_ms, 0),
    }

    # Extract structured tables/charts from the PDF
    if extract_charts:
        try:
            charts = _extract_charts_from_pdf(file_path)
            if charts:
                result["charts"] = charts
                chart_text = _format_charts(charts)
                result["text"] = result["text"] + "\n\n" + chart_text
        except Exception as exc:
            logger.warning(f"Chart extraction failed for PDF: {exc}")

    return result


def _pdf_ocr(file_path: Path) -> str:
    """OCR a PDF using pdftoppm + tesseract.

    Searches for TESSDATA_PREFIX in environment variables first,
    then common system paths, then ~/.local/share/tessdata.
    Supports chi_sim+eng language pack (auto-detect Chinese characters).
    """
    import os as _os

    # Resolve tessdata prefix — needed for custom installs where
    # tesseract lang data is not in the default /usr/share path.
    _tessdata_prefix = _os.environ.get("TESSDATA_PREFIX", "")
    if not _tessdata_prefix:
        for _candidate in (
            "/usr/share/tesseract-ocr/4.00",
            "/usr/share/tesseract-ocr",
            str(Path.home() / ".local" / "share" / "tessdata"),
        ):
            if Path(_candidate, "tessdata", "eng.traineddata").exists() or \
               Path(_candidate, "eng.traineddata").exists():
                _tessdata_prefix = _candidate
                break

    _env = dict(_os.environ)
    if _tessdata_prefix:
        _env["TESSDATA_PREFIX"] = _tessdata_prefix
        logger.info(f"OCR: TESSDATA_PREFIX={_tessdata_prefix}")

    try:
        # Convert PDF pages to images
        pages_dir = tempfile.mkdtemp(prefix="miqi_ocr_")
        try:
            subprocess.run(
                ["pdftoppm", "-r", "200", "-png", str(file_path), f"{pages_dir}/page"],
                capture_output=True, timeout=120,
            )

            page_images = sorted(Path(pages_dir).glob("page-*.png"))
            if not page_images:
                logger.warning("pdftoppm produced no images")
                return ""

            # OCR each page
            full_text_parts = []
            for img_path in page_images:
                try:
                    result = subprocess.run(
                        ["tesseract", str(img_path), "stdout", "-l", "chi_sim+eng"],
                        capture_output=True, text=True, timeout=30,
                        env=_env,
                    )
                    if result.returncode == 0 and result.stdout.strip():
                        full_text_parts.append(result.stdout.strip())
                except Exception as exc:
                    logger.warning(f"Tesseract OCR failed for {img_path}: {exc}")

            return "\n\n".join(full_text_parts)
        finally:
            # Guarantee cleanup on every exit path (including no-images + OCR failures)
            shutil.rmtree(pages_dir, ignore_errors=True)
    except FileNotFoundError:
        logger.warning("pdftoppm or tesseract not found, OCR unavailable")
        return ""
    except subprocess.TimeoutExpired:
        logger.warning("PDF OCR timed out (too many pages or complex layout)")
        return ""
    except Exception as exc:
        logger.error(f"OCR pipeline failed: {exc}")
        return ""


# ── Images (OCR) ─────────────────────────────────────────────────────────

def _image_ocr(file_path: Path) -> str:
    """OCR a single image.

    Prefers RapidOCR (pure-Python ONNX runtime, no system binary, strong
    zh/en accuracy) and falls back to tesseract (chi_sim+eng) when RapidOCR
    is not installed. Returns "" (never raises) when no engine is available.
    """
    # Engine 1: RapidOCR — onnxruntime, no system dependency (#659).
    try:
        from rapidocr_onnxruntime import RapidOCR

        _engine = RapidOCR()
        result, _ = _engine(str(file_path))
        if result:
            lines = [line[1] for line in result if len(line) > 1 and line[1].strip()]
            if lines:
                return "\n".join(lines)
    except Exception as exc:
        logger.warning(f"RapidOCR unavailable for {file_path}: {exc}")

    # Engine 2: tesseract (chi_sim+eng) — needs the system binary.
    _tessdata_prefix = os.environ.get("TESSDATA_PREFIX", "")
    if not _tessdata_prefix:
        for _candidate in (
            "/usr/share/tesseract-ocr/4.00",
            "/usr/share/tesseract-ocr",
            str(Path.home() / ".local" / "share" / "tessdata"),
        ):
            if Path(_candidate, "tessdata", "eng.traineddata").exists() or \
               Path(_candidate, "eng.traineddata").exists():
                _tessdata_prefix = _candidate
                break

    _env = dict(os.environ)
    if _tessdata_prefix:
        _env["TESSDATA_PREFIX"] = _tessdata_prefix

    try:
        result = subprocess.run(
            ["tesseract", str(file_path), "stdout", "-l", "chi_sim+eng"],
            capture_output=True, text=True, timeout=60, env=_env,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except FileNotFoundError:
        logger.warning("tesseract not found, image OCR unavailable")
    except subprocess.TimeoutExpired:
        logger.warning("image OCR timed out")
    except Exception as exc:
        logger.warning(f"image OCR failed: {exc}")
    return ""


def _parse_image(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Parse an image: OCR text (tesseract, chi_sim+eng) + basic metadata.

    Falls back to metadata-only when tesseract is unavailable. The text is
    prefixed with an image header (size/format) so the model understands it
    is describing an image.
    """
    result: dict[str, Any] = {
        "text": "", "page_count": 1, "size_bytes": file_path.stat().st_size,
        "mime_type": _SUFFIX_TO_MIME.get(file_path.suffix.lower(), "image/octet-stream"),
        "ocr_used": False, "parse_ms": 0, "charts": [],
    }
    _t0 = time.time()
    try:
        from PIL import Image as _PILImage
        with _PILImage.open(file_path) as img:
            width, height = img.size
            fmt = img.format or ""
        header = f"[图片] {file_path.name} ({width}x{height}, {fmt})\n"
        ocr_text = _image_ocr(file_path)
        if ocr_text:
            result["ocr_used"] = True
            result["text"] = (header + ocr_text)[:max_chars]
        else:
            result["text"] = (header + "（图片未提取到文字：tesseract 不可用或图中无文字）")[:max_chars]
        result["parse_ms"] = int((time.time() - _t0) * 1000)
    except Exception as exc:
        logger.warning(f"image parse failed for {file_path}: {exc}")
        result["text"] = (f"[图片] {file_path.name}（解析失败: {exc}）")[:max_chars]
    return result


# ── Chart & Table extraction ────────────────────────────────────────────────

def _extract_chart_data_from_pdf_page(page_obj: Any) -> list[dict[str, Any]]:
    """Attempt to extract structured chart/table data from a PDF page.

    Uses pdfplumber when available for robust table extraction.
    Returns a list of table dicts with headers, rows, and caption info.
    """
    tables = []
    try:
        import pdfplumber
        # pdfplumber works with file paths, not page objects
    except ImportError:
        return tables
    return tables


def _extract_charts_from_pdf(file_path: Path) -> list[dict[str, Any]]:
    """Extract structured tables and chart data from a PDF using pdfplumber.

    Returns list of {page, table_index, headers, rows, caption}.
    """
    charts = []
    try:
        import pdfplumber
        with pdfplumber.open(str(file_path)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_tables = page.extract_tables()
                for t_idx, table in enumerate(page_tables):
                    if not table or len(table) < 2:
                        continue
                    # First row as headers, rest as data
                    headers = [str(h) if h else "" for h in table[0]]
                    rows = [
                        [str(c) if c else "" for c in row]
                        for row in table[1:]
                    ]
                    if any(any(c for c in row) for row in rows):
                        charts.append({
                            "page": page_num,
                            "table_index": t_idx,
                            "headers": headers,
                            "rows": rows,
                            "caption": f"Table {t_idx + 1} on page {page_num}",
                        })
    except ImportError:
        pass
    except Exception as exc:
        logger.warning(f"pdfplumber chart extraction failed: {exc}")
    return charts


def _format_charts(charts: list[dict[str, Any]]) -> str:
    """Format extracted chart/table data as readable text for LLM."""
    if not charts:
        return ""
    parts = ["\n## Extracted Tables & Charts"]
    for chart in charts:
        parts.append(f"\n### {chart['caption']}")
        parts.append(" | ".join(chart["headers"]))
        parts.append("-" * 40)
        for row in chart["rows"][:50]:  # Limit rows
            parts.append(" | ".join(row))
    return "\n".join(parts)


# ── Word (DOCX) Parsing ────────────────────────────────────────────────────

def _parse_docx(file_path: Path, *, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Extract text from a Word document."""
    import time
    t0 = time.monotonic()

    text = ""
    page_count = 1

    try:
        from docx import Document
        doc = Document(str(file_path))
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                # Detect headings
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    paragraphs.append(f"## {para.text}")
                else:
                    paragraphs.append(para.text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                paragraphs.append(" | ".join(cells))

        text = "\n\n".join(paragraphs)
        page_count = max(1, len(doc.paragraphs) // 40)
    except ImportError:
        logger.warning("python-docx not installed")
        text = "[文档解析需要 python-docx 库]"
    except Exception as exc:
        logger.error(f"DOCX parsing failed: {exc}")
        text = f"[文档解析失败: {exc}]"

    text = text[:max_chars]
    parse_ms = (time.monotonic() - t0) * 1000

    return {
        "text": text,
        "page_count": page_count,
        "size_bytes": file_path.stat().st_size,
        "mime_type": _SUFFIX_TO_MIME.get(file_path.suffix.lower(), "application/octet-stream"),
        "ocr_used": False,
        "parse_ms": round(parse_ms, 0),
    }


# ── PowerPoint (PPTX) Parsing ──────────────────────────────────────────────

def _parse_pptx(file_path: Path, *, max_chars: int = MAX_CONTEXT_CHARS,
                extract_charts: bool = True) -> dict[str, Any]:
    """Extract text from a PowerPoint presentation."""
    import time
    t0 = time.monotonic()

    text = ""
    slide_count = 0

    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            lines = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            lines.append(para.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        lines.append(" | ".join(cells))
            slides_text.append("\n".join(lines))
            slide_count = i

        # Also extract notes
        for i, slide in enumerate(prs.slides, 1):
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slides_text[i - 1] += f"\n[Notes]: {notes}"

        text = "\n\n".join(slides_text)
    except ImportError:
        logger.warning("python-pptx not installed")
        text = "[文档解析需要 python-pptx 库]"
    except Exception as exc:
        logger.error(f"PPTX parsing failed: {exc}")
        text = f"[文档解析失败: {exc}]"

    text = text[:max_chars]
    parse_ms = (time.monotonic() - t0) * 1000

    # Extract charts/tables from PPTX if requested
    charts = []
    if extract_charts and slide_count > 0:
        try:
            import pptx as _pptx_check
            prs2 = _pptx_check.Presentation(str(file_path))
            for i, slide in enumerate(prs2.slides, 1):
                for shape in slide.shapes:
                    if shape.has_chart:
                        try:
                            ch = shape.chart
                            chart_info = {
                                "slide": i,
                                "chart_type": str(ch.chart_type) if hasattr(ch, 'chart_type') else "unknown",
                                "has_title": ch.has_title,
                            }
                            if ch.has_title:
                                chart_info["title"] = ch.chart_title.text_frame.text if ch.chart_title.has_text_frame else ""
                            # Try to get series data
                            series_data = []
                            for s_idx, series in enumerate(ch.series):
                                try:
                                    vals = list(series.values)
                                    series_data.append({
                                        "name": str(series.format.fill) if hasattr(series, 'format') else f"Series {s_idx}",
                                        "values": [str(v) for v in vals[:20]],
                                    })
                                except Exception:
                                    pass
                            if series_data:
                                chart_info["series"] = series_data
                            charts.append(chart_info)
                        except Exception:
                            pass
        except ImportError:
            pass
        except Exception as exc:
            logger.warning(f"PPTX chart extraction failed: {exc}")

    if charts:
        chart_text = "\n## Charts in Presentation\n"
        for c in charts:
            chart_text += f"\n### Slide {c['slide']}: {c.get('title', c.get('chart_type', 'Chart'))}\n"
            for s in c.get("series", []):
                chart_text += f"- {s['name']}: {', '.join(s.get('values', [])[:10])}\n"
        text = text + chart_text

    return {
        "text": text,
        "page_count": slide_count,
        "size_bytes": file_path.stat().st_size,
        "mime_type": _SUFFIX_TO_MIME.get(file_path.suffix.lower(), "application/octet-stream"),
        "ocr_used": False,
        "parse_ms": round(parse_ms, 0),
    }


# ── Excel (XLSX) Parsing ───────────────────────────────────────────────────

def _parse_xlsx(file_path: Path, *, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Extract text and data from an Excel spreadsheet."""
    import time
    t0 = time.monotonic()

    text = ""
    sheet_count = 0

    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(file_path), data_only=True, read_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_count += 1
            lines = [f"=== Sheet: {sheet_name} ==="]

            # Collect rows
            row_count = 0
            max_rows = 500  # Limit rows to avoid huge output
            for row in ws.iter_rows(values_only=True):
                row_count += 1
                if row_count > max_rows:
                    remaining = ws.max_row - row_count if ws.max_row else 0
                    lines.append(f"... ({remaining} more rows)")
                    break
                cells = [str(cell) if cell is not None else "" for cell in row]
                lines.append(" | ".join(cells))

            sheets_text.append("\n".join(lines))

        text = "\n\n".join(sheets_text)
        wb.close()
    except ImportError:
        logger.warning("openpyxl not installed")
        text = "[文档解析需要 openpyxl 库]"
    except Exception as exc:
        logger.error(f"XLSX parsing failed: {exc}")
        text = f"[文档解析失败: {exc}]"

    text = text[:max_chars]
    parse_ms = (time.monotonic() - t0) * 1000

    return {
        "text": text,
        "page_count": sheet_count,
        "size_bytes": file_path.stat().st_size,
        "mime_type": _SUFFIX_TO_MIME.get(file_path.suffix.lower(), "application/octet-stream"),
        "ocr_used": False,
        "parse_ms": round(parse_ms, 0),
    }


# ── Markdown Parsing ────────────────────────────────────────────────────────

_MERMAID_RE = re.compile(r"```mermaid\s*\n(.*?)```", re.DOTALL | re.IGNORECASE)
_CODE_BLOCK_RE = re.compile(r"```(\w+)?\s*\n(.*?)```", re.DOTALL)
_TABLE_SEP_RE = re.compile(r"^\s*\|?[\s\-:]+\|[\s\-:|]+\s*$")
_TABLE_ROW_RE = re.compile(r"^\s*\|.+\|\s*$")


def _parse_markdown(file_path: Path, *, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read a Markdown file with structured extraction.

    Extracts:
    - Full raw text for LLM context
    - Mermaid diagrams with their code blocks
    - Markdown tables (pipe-style) with headers+rows
    - Code blocks with language detection
    - Heading outline for document structure
    """
    import time
    t0 = time.monotonic()

    try:
        raw = file_path.read_text(encoding="utf-8")
    except Exception as exc:
        logger.error(f"Markdown read failed: {exc}")
        return {
            "text": f"[文件读取失败: {exc}]",
            "page_count": 1,
            "size_bytes": 0,
            "mime_type": "text/markdown",
            "ocr_used": False,
            "parse_ms": 0,
        }

    size_bytes = file_path.stat().st_size

    # Extract mermaid diagrams
    mermaids = []
    for m in _MERMAID_RE.finditer(raw):
        code = m.group(1).strip()
        chart_type = "mermaid"
        first_line = code.split("\n")[0].strip() if code else ""
        if first_line in ("graph", "flowchart", "sequenceDiagram", "classDiagram",
                          "stateDiagram", "erDiagram", "gantt", "pie", "gitGraph", "mindmap"):
            chart_type = first_line
        mermaids.append({"type": chart_type, "code": code})

    # Extract code blocks (excluding mermaid)
    code_blocks = []
    for m in _CODE_BLOCK_RE.finditer(raw):
        lang = (m.group(1) or "").lower()
        if lang == "mermaid":
            continue
        code = m.group(2).strip()
        code_blocks.append({"language": lang or "text", "code": code[:5000]})

    # Extract pipe tables
    tables = []
    lines = raw.split("\n")
    i = 0
    while i < len(lines):
        if i + 1 < len(lines) and _TABLE_ROW_RE.match(lines[i]) and _TABLE_SEP_RE.match(lines[i + 1]):
            headers = [h.strip() for h in lines[i].strip("|").split("|")]
            rows = []
            j = i + 2
            while j < len(lines) and _TABLE_ROW_RE.match(lines[j]):
                cells = [c.strip() for c in lines[j].strip("|").split("|")]
                while len(cells) < len(headers):
                    cells.append("")
                rows.append(cells[:len(headers)])
                j += 1
            if rows:
                tables.append({"headers": headers, "rows": rows, "line": i + 1})
            i = j
        else:
            i += 1

    # Extract heading outline
    heading_re = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)
    headings = [{"level": len(m.group(1)), "text": m.group(2).strip()}
                for m in heading_re.finditer(raw)]

    # Build structured text for LLM
    text_parts = [raw]

    if mermaids:
        text_parts.append("\n\n## Mermaid Diagrams")
        for idx, d in enumerate(mermaids, 1):
            text_parts.append(f"\n### Diagram {idx} ({d['type']})\n```mermaid\n{d['code']}\n```")

    if tables:
        text_parts.append("\n\n## Tables")
        for idx, t in enumerate(tables, 1):
            text_parts.append(f"\n### Table {idx}")
            text_parts.append(" | ".join(t["headers"]))
            text_parts.append("-" * 40)
            for row in t["rows"]:
                text_parts.append(" | ".join(row))

    text = "\n".join(text_parts)[:max_chars]
    parse_ms = (time.monotonic() - t0) * 1000

    result = {
        "text": text,
        "page_count": 1,
        "size_bytes": size_bytes,
        "mime_type": "text/markdown",
        "ocr_used": False,
        "parse_ms": round(parse_ms, 0),
    }

    if headings:
        result["headings"] = headings[:50]
    if mermaids:
        result["mermaids"] = mermaids
    if tables:
        result["tables"] = tables
    if code_blocks:
        result["code_blocks"] = code_blocks[:20]

    return result


# ── HTML Parser ──────────────────────────────────────────────────────────

def _parse_html(file_path: Path, max_chars: int = 50000) -> dict:
    """Extract text from HTML files using lxml."""
    if not _HAS_LXML_HTML:
        raise RuntimeError("lxml is required for HTML parsing")

    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")

    try:
        doc = _lxml_html.document_fromstring(raw)
    except Exception:
        # Fallback: plain text stripping via stdlib
        stem = Path(file_path.stem).stem if file_path.stem else "HTML"
        from html.parser import HTMLParser as _StdlibParser
        class _Stripper(_StdlibParser):
            def __init__(self):
                super().__init__()
                self.text: list[str] = []
            def handle_data(self, data):
                self.text.append(data)
        s = _Stripper()
        s.feed(raw)
        text = " ".join(s.text)
        text = re.sub(r"\s+", " ", text).strip()
        return {
            "text": text[:max_chars],
            "page_count": 1,
            "size_bytes": file_path.stat().st_size,
            "mime_type": "text/html",
            "ocr_used": False,
            "parse_ms": round((time.perf_counter() - t0) * 1000),
        }

    # Extract title BEFORE removing head (the title lives inside <head>)
    title_el = doc.xpath("//title/text()")
    title = title_el[0].strip() if title_el else ""

    # Remove script/style/noscript tags and head/meta/link elements
    # Use drop_tree() — safer than parent.remove() which can fail if parent is None
    for tag in doc.xpath("//script|//style|//noscript|//head|//meta|//link"):
        tag.drop_tree()

    # Get visible text from body
    body = doc.xpath("//body")
    body_text = " ".join(body[0].itertext()) if body else ""
    body_text = re.sub(r"\s+", " ", body_text).strip()

    parts = []
    if title:
        parts.append(f"Title: {title}")
    if body_text:
        parts.append(body_text)

    text = "\n\n".join(parts)
    return {
        "text": text[:max_chars],
        "page_count": 1,
        "size_bytes": file_path.stat().st_size,
        "mime_type": "text/html",
        "ocr_used": False,
        "parse_ms": round((time.perf_counter() - t0) * 1000),
    }


# ── CSV Parser ────────────────────────────────────────────────────────────

def _parse_csv(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Extract and format CSV data as readable text."""
    import csv as _csv_mod
    t0 = time.perf_counter()
    try:
        raw = file_path.read_text(encoding="utf-8", errors="replace")
        reader = _csv_mod.reader(io.StringIO(raw, newline=""))
        rows = list(reader)
        if not rows:
            return _make_text_result(file_path, raw, max_chars, t0, "text/csv")
        # Format as table-like text
        lines = []
        headers = rows[0] if rows else []
        if headers:
            lines.append(" | ".join(str(h) for h in headers))
            lines.append("-" * 40)
        for row in rows[1:101]:  # max 100 data rows
            lines.append(" | ".join(str(c) for c in row))
        if len(rows) > 101:
            lines.append(f"... ({len(rows) - 101} more rows)")
        text = "\n".join(lines)
    except Exception as exc:
        logger.warning(f"CSV parsing failed: {exc}, falling back to raw text")
        raw = file_path.read_text(encoding="utf-8", errors="replace")
        text = raw
    return _make_text_result(file_path, text, max_chars, t0, "text/csv")


# ── JSON Parser ───────────────────────────────────────────────────────────

def _parse_json(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Extract and format JSON data as readable text."""
    import json as _json_mod
    t0 = time.perf_counter()
    try:
        raw = file_path.read_text(encoding="utf-8", errors="replace")
        data = _json_mod.loads(raw)
        text = _json_mod.dumps(data, ensure_ascii=False, indent=2)
    except Exception as exc:
        logger.warning(f"JSON parsing failed: {exc}, falling back to raw text")
        text = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, text, max_chars, t0, "application/json")


# ── XML File Parser ───────────────────────────────────────────────────────

def _parse_xml_file(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Extract text from XML files via lxml or plain-text fallback."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    if _HAS_LXML_HTML:
        try:
            doc = _lxml_html.document_fromstring(raw)
            text = " ".join(doc.itertext())
            text = re.sub(r"\s+", " ", text).strip()
            return _make_text_result(file_path, text, max_chars, t0, "application/xml")
        except Exception as exc:
            logger.warning(f"lxml XML parsing failed: {exc}, falling back to raw text")
    return _make_text_result(file_path, raw, max_chars, t0, "application/xml")


# ── YAML Parser ───────────────────────────────────────────────────────────

def _parse_yaml(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read YAML file as plain text (structure preservation)."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, raw, max_chars, t0, "text/yaml")


# ── ENV Parser ────────────────────────────────────────────────────────────

def _parse_env(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read .env file as plain text."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, raw, max_chars, t0, "text/plain")


# ── LOG Parser ────────────────────────────────────────────────────────────

def _parse_log(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read log file as plain text."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, raw, max_chars, t0, "text/plain")


# ── SQL Parser ────────────────────────────────────────────────────────────

def _parse_sql(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read SQL file as plain text."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, raw, max_chars, t0, "text/x-sql")


# ── INI Parser ────────────────────────────────────────────────────────────

def _parse_ini(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read INI file as plain text."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, raw, max_chars, t0, "text/plain")


# ── TOML Parser ───────────────────────────────────────────────────────────

def _parse_toml(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read TOML file as plain text."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, raw, max_chars, t0, "text/plain")


# ── HTACCESS Parser ───────────────────────────────────────────────────────

def _parse_htaccess(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read .htaccess file as plain text."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, raw, max_chars, t0, "text/plain")


# ── Shell Script Parser ───────────────────────────────────────────────────

def _parse_sh(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read shell script as plain text."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, raw, max_chars, t0, "text/x-shellscript")


# ── TXT Parser ────────────────────────────────────────────────────────────

def _parse_txt(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Read plain text file."""
    t0 = time.perf_counter()
    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, raw, max_chars, t0, "text/plain")


# ── RTF Parser ────────────────────────────────────────────────────────────

def _parse_rtf(file_path: Path, max_chars: int = MAX_CONTEXT_CHARS) -> dict[str, Any]:
    """Extract text from RTF files by stripping RTF markup."""
    t0 = time.perf_counter()
    try:
        raw = file_path.read_text(encoding="utf-8", errors="replace")
        # Decode RTF Unicode escapes (\uN with ANSI fallback) before stripping markup.
        # \u233?  → é (U+00E9, decimal 233, fallback ?)
        text = re.sub(
            r"\\u(-?\d+)\s*\??",
            lambda m: chr(int(m.group(1)) % 0x10000) if 0 < int(m.group(1)) <= 0x10FFFF else m.group(0),
            raw,
        )
        # Decode \'xx hex byte escapes (code-page dependent; basic ASCII/Latin-1 pass-through)
        text = re.sub(r"\\\'([0-9a-fA-F]{2})", lambda m: chr(int(m.group(1), 16)), text)
        # Strip remaining RTF control words (rtf1, b, i, par, etc.)
        text = re.sub(r"\\[a-zA-Z]+\s?", " ", text)
        # Remove { and }
        text = text.replace("{", "").replace("}", "")
        # Normalize whitespace
        text = re.sub(r"\s+", " ", text).strip()
        if len(text.strip()) < 5:
            # Fallback: RTF parsing produced too little text, return raw
            text = raw
    except Exception as exc:
        logger.warning(f"RTF parsing failed: {exc}")
        text = file_path.read_text(encoding="utf-8", errors="replace")
    return _make_text_result(file_path, text, max_chars, t0, "text/rtf")


# ── Shared helpers ────────────────────────────────────────────────────────

def _make_text_result(
    file_path: Path,
    text: str,
    max_chars: int,
    t0: float,
    mime_type: str = "text/plain",
) -> dict[str, Any]:
    """Build a standard result dict for text-based parsers."""
    return {
        "text": text[:max_chars],
        "page_count": 1,
        "size_bytes": file_path.stat().st_size,
        "mime_type": mime_type,
        "ocr_used": False,
        "parse_ms": round((time.perf_counter() - t0) * 1000),
    }
