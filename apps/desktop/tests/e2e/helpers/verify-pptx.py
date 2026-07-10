"""Verify latest pptx-generator output in workspace matches prompt spec."""
import json
import os
import sys
import glob as _glob
from pathlib import Path
from pptx import Presentation

# Force UTF-8 stdout so Unicode bullets/symbols don't crash on GBK Windows
sys.stdout.reconfigure(encoding='utf-8')

workspace = sys.argv[1]
files = _glob.glob(os.path.join(workspace, "**", "*.pptx"), recursive=True)
if not files:
    json.dump({"pass": False, "checks": [{"label":"find pptx","pass":False,"detail":"no pptx found"}]}, sys.stdout)
    sys.exit(1)

filepath = max(files, key=os.path.getmtime)
prs = Presentation(filepath)

texts = []
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    texts.append(t)

all_text = "\n".join(texts)
result = {
    "slides": len(prs.slides),
    "texts": texts,
    "pass": True,
    "checks": [],
}

def check(label, condition, detail=""):
    result["checks"].append({"label": label, "pass": bool(condition), "detail": detail})
    if not condition:
        result["pass"] = False

check("slide count >= 5", len(prs.slides) >= 5, f"got {len(prs.slides)}")
check("cover title", texts[0] == "人工智能简介")
check("cover subtitle", texts[1] == "技术、应用与未来")
for kw in ["什么是AI", "核心技术", "应用场景", "未来展望"]:
    check(f"TOC: {kw}", kw in all_text)
for kw in ["机器学习", "深度学习", "NLP"]:
    check(f"content: {kw}", kw in all_text)
summary_kws = ["AI重塑行业", "人机协作", "安全对齐", "拥抱AI"]
summary_found = sum(1 for kw in summary_kws if kw in all_text)
check(f"summary keywords >= 2 (found {summary_found})", summary_found >= 2)

json.dump(result, sys.stdout, ensure_ascii=False, indent=2)
sys.exit(0 if result["pass"] else 1)
