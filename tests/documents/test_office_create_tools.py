"""Tests for Office document creation tools."""

import pytest


def _east_asia_font(run):
    from docx.oxml.ns import qn

    r_fonts = run._element.rPr.rFonts
    return r_fonts.get(qn("w:eastAsia")) if r_fonts is not None else None


@pytest.mark.asyncio
async def test_create_docx_supports_title_paragraphs_and_tables(tmp_path):
    from docx import Document

    from miqi.documents.docx_tool import CreateDocxTool

    files_dir = tmp_path / "files"
    tool = CreateDocxTool(workspace=files_dir, allowed_dir=files_dir)

    result = await tool.execute(
        filename="report",
        title="Quarterly Report",
        paragraphs=["Summary paragraph."],
        tables=[{"rows": [["Metric", "Value"], ["Revenue", 42]]}],
    )

    path = files_dir / "report.docx"
    assert "Created:" in result
    assert path.exists()

    doc = Document(str(path))
    texts = [p.text for p in doc.paragraphs if p.text.strip()]
    assert "Quarterly Report" in texts
    assert "Summary paragraph." in texts
    assert doc.tables[0].cell(1, 0).text == "Revenue"


@pytest.mark.asyncio
async def test_create_docx_parses_markdown_heading_levels(tmp_path):
    from docx import Document

    from miqi.documents.docx_tool import CreateDocxTool

    files_dir = tmp_path / "files"
    tool = CreateDocxTool(workspace=files_dir, allowed_dir=files_dir)

    result = await tool.execute(
        filename="headings",
        content="# Level 1\n### Level 3\n正文",
    )

    path = files_dir / "headings.docx"
    assert "Created:" in result

    doc = Document(str(path))
    paragraphs = [p for p in doc.paragraphs if p.text.strip()]
    assert [p.text for p in paragraphs] == ["Level 1", "Level 3", "正文"]
    assert paragraphs[0].style.name == "Heading 1"
    assert paragraphs[1].style.name == "Heading 3"


@pytest.mark.asyncio
async def test_docx_read_resolves_workspace_relative_paths(tmp_path):
    from miqi.documents.docx_tool import CreateDocxTool, DocxReadTool

    files_dir = tmp_path / "files"
    create = CreateDocxTool(workspace=files_dir, allowed_dir=files_dir)
    await create.execute(filename="relative_doc", title="Relative Title")

    read = DocxReadTool(workspace=files_dir, allowed_dir=files_dir)
    result = await read.execute(filename="relative_doc.docx")

    assert "Relative Title" in result


@pytest.mark.asyncio
async def test_create_docx_applies_chinese_title_and_body_formatting(tmp_path):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    from miqi.documents.docx_tool import CreateDocxTool

    files_dir = tmp_path / "files"
    tool = CreateDocxTool(workspace=files_dir, allowed_dir=files_dir)

    result = await tool.execute(
        filename="essay",
        title="那山，那水，那故乡",
        paragraphs=["小时候，故乡在我眼中不过是一个地理名词。"],
        style_preset="chinese_document",
    )

    path = files_dir / "essay.docx"
    assert "Created:" in result

    doc = Document(str(path))
    title = next(p for p in doc.paragraphs if p.text == "那山，那水，那故乡")
    body = next(p for p in doc.paragraphs if p.text.startswith("小时候"))

    assert title.alignment == WD_ALIGN_PARAGRAPH.CENTER
    assert title.runs[0].bold is True
    assert title.runs[0].font.size == Pt(16)
    assert _east_asia_font(title.runs[0]) == "黑体"
    assert body.runs[0].font.size == Pt(12)
    assert _east_asia_font(body.runs[0]) == "宋体"
    assert body.paragraph_format.line_spacing == 1.5


@pytest.mark.asyncio
async def test_edit_docx_can_apply_formatting_without_text_change(tmp_path):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    from miqi.documents.docx_tool import EditDocxTool

    files_dir = tmp_path / "files"
    files_dir.mkdir()
    path = files_dir / "essay.docx"
    doc = Document()
    doc.add_heading("那山，那水，那故乡", level=0)
    doc.add_paragraph("小时候，故乡在我眼中不过是一个地理名词。")
    doc.save(str(path))

    tool = EditDocxTool(workspace=files_dir, allowed_dir=files_dir)
    result = await tool.execute(
        filename="essay.docx",
        format_instructions="正文宋体小四，段落1.5行距，标题黑体加粗，三号字居中",
    )

    assert "formatted" in result
    edited = Document(str(path))
    title = next(p for p in edited.paragraphs if p.text == "那山，那水，那故乡")
    body = next(p for p in edited.paragraphs if p.text.startswith("小时候"))

    assert title.alignment == WD_ALIGN_PARAGRAPH.CENTER
    assert title.runs[0].bold is True
    assert title.runs[0].font.size == Pt(16)
    assert _east_asia_font(title.runs[0]) == "黑体"
    assert body.runs[0].font.size == Pt(12)
    assert _east_asia_font(body.runs[0]) == "宋体"
    assert body.paragraph_format.line_spacing == 1.5


@pytest.mark.asyncio
async def test_edit_docx_format_instructions_override_bad_structured_style(tmp_path):
    from docx import Document
    from docx.shared import Pt

    from miqi.documents.docx_tool import EditDocxTool

    files_dir = tmp_path / "files"
    files_dir.mkdir()
    path = files_dir / "essay.docx"
    doc = Document()
    doc.add_heading("那山，那水，那故乡", level=1)
    doc.add_paragraph("小时候，故乡在我眼中不过是一个地理名词。")
    doc.save(str(path))

    tool = EditDocxTool(workspace=files_dir, allowed_dir=files_dir)
    result = await tool.execute(
        filename="essay.docx",
        title_style={"font_size_pt": 14},
        body_style={"font_size_pt": 16},
        format_instructions="正文宋体小四，段落1.5行距，标题黑体加粗，三号字居中",
    )

    assert "formatted" in result
    edited = Document(str(path))
    title = next(p for p in edited.paragraphs if p.text == "那山，那水，那故乡")
    body = next(p for p in edited.paragraphs if p.text.startswith("小时候"))

    assert title.runs[0].font.size == Pt(16)
    assert body.runs[0].font.size == Pt(12)


@pytest.mark.asyncio
async def test_create_xlsx_supports_multiple_sheets_formulas_and_charts(tmp_path):
    from openpyxl import load_workbook

    from miqi.documents.xlsx_tool import CreateXlsxTool

    files_dir = tmp_path / "files"
    tool = CreateXlsxTool(workspace=files_dir, allowed_dir=files_dir)

    result = await tool.execute(
        filename="analysis",
        sheets=[
            {
                "name": "Sales",
                "rows": [["Month", "Sales"], ["Jan", 10], ["Feb", 20], ["Total", "=SUM(B2:B3)"]],
                "charts": [
                    {
                        "type": "bar",
                        "title": "Sales",
                        "data_range": "B1:B3",
                        "category_range": "A2:A3",
                        "anchor": "D2",
                    }
                ],
            },
            {"name": "Notes", "rows": [["Ready"]]},
        ],
    )

    path = files_dir / "analysis.xlsx"
    assert "Created:" in result
    assert path.exists()

    wb = load_workbook(str(path), data_only=False)
    assert wb.sheetnames == ["Sales", "Notes"]
    assert wb["Sales"]["B4"].value == "=SUM(B2:B3)"
    assert len(wb["Sales"]._charts) == 1


@pytest.mark.asyncio
async def test_create_xlsx_accepts_series_based_chart_specs(tmp_path):
    from openpyxl import load_workbook

    from miqi.documents.xlsx_tool import CreateXlsxTool

    files_dir = tmp_path / "files"
    tool = CreateXlsxTool(workspace=files_dir, allowed_dir=files_dir)

    result = await tool.execute(
        filename="sales_analysis",
        sheets=[
            {
                "name": "Sales",
                "rows": [
                    ["Month", "Product A", "Product B", "Total"],
                    ["Jan", 100, 80, "=SUM(B2:C2)"],
                    ["Feb", 120, 90, "=SUM(B3:C3)"],
                    ["Mar", 150, 110, "=SUM(B4:C4)"],
                    ["Apr", 130, 140, "=SUM(B5:C5)"],
                ],
                "charts": [
                    {
                        "type": "bar",
                        "title": "Monthly Sales",
                        "categories": ["Jan", "Feb", "Mar", "Apr"],
                        "series": [
                            {"name": "Product A", "values": [100, 120, 150, 130]},
                            {"name": "Product B", "values": [80, 90, 110, 140]},
                        ],
                        "anchor_cell": "F2",
                    }
                ],
            }
        ],
    )

    path = files_dir / "sales_analysis.xlsx"
    assert "Created:" in result
    wb = load_workbook(str(path), data_only=False)
    assert len(wb["Sales"]._charts) == 1


@pytest.mark.asyncio
async def test_create_xlsx_accepts_top_level_rows_and_rejects_invalid_sheets(tmp_path):
    from openpyxl import load_workbook

    from miqi.documents.xlsx_tool import CreateXlsxTool, XlsxReadTool

    files_dir = tmp_path / "files"
    tool = CreateXlsxTool(workspace=files_dir, allowed_dir=files_dir)

    result = await tool.execute(
        filename="top_rows",
        sheet_name="Data",
        rows=[["A", "B"], [1, 2]],
    )

    path = files_dir / "top_rows.xlsx"
    assert "Created:" in result
    wb = load_workbook(str(path), data_only=False)
    assert wb.sheetnames == ["Data"]
    assert wb["Data"]["B2"].value == 2

    read = XlsxReadTool(workspace=files_dir, allowed_dir=files_dir)
    read_result = await read.execute(filename="top_rows.xlsx", sheet_name="Data")
    assert "1 | 2" in read_result

    invalid = await tool.execute(filename="bad", sheets="not valid", rows=[["A"]])
    assert "sheets must be an object or array" in invalid


@pytest.mark.asyncio
async def test_create_pptx_supports_multiple_slides_and_bullets(tmp_path):
    from pptx import Presentation

    from miqi.documents.pptx_tool import CreatePptxTool

    files_dir = tmp_path / "files"
    tool = CreatePptxTool(workspace=files_dir, allowed_dir=files_dir)

    result = await tool.execute(
        filename="deck",
        slides=[
            {"title": "Overview", "content": "Opening"},
            {"title": "Plan", "bullets": ["Build", "Verify", "Ship"]},
        ],
    )

    path = files_dir / "deck.pptx"
    assert "Created:" in result
    assert path.exists()

    prs = Presentation(str(path))
    assert len(prs.slides) == 2
    text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Overview" in text
    assert "Build" in text


@pytest.mark.asyncio
async def test_create_pptx_supports_subtitle_array_content_and_relative_read(tmp_path):
    from pptx import Presentation

    from miqi.documents.pptx_tool import CreatePptxTool, PptxReadTool

    files_dir = tmp_path / "files"
    tool = CreatePptxTool(workspace=files_dir, allowed_dir=files_dir)

    result = await tool.execute(
        filename="deck",
        slides=[
            {"title": "Cover", "subtitle": "Sub title"},
            {"title": "Agenda", "content": ["A", "B"]},
        ],
    )

    path = files_dir / "deck.pptx"
    assert "Created:" in result
    prs = Presentation(str(path))
    text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Sub title" in text
    assert "A" in text
    assert "B" in text
    assert "['A', 'B']" not in text

    read = PptxReadTool(workspace=files_dir, allowed_dir=files_dir)
    read_result = await read.execute(filename="deck.pptx")
    assert "Presentation: 2 slides" in read_result


@pytest.mark.asyncio
async def test_edit_docx_replaces_and_appends_text(tmp_path):
    from docx import Document

    from miqi.documents.docx_tool import EditDocxTool

    files_dir = tmp_path / "files"
    files_dir.mkdir()
    path = files_dir / "report.docx"
    doc = Document()
    doc.add_paragraph("Draft status")
    doc.save(str(path))

    tool = EditDocxTool(workspace=files_dir, allowed_dir=files_dir)
    result = await tool.execute(
        filename="report.docx",
        old_text="Draft",
        new_text="Final",
        append_paragraphs=["Approved"],
    )

    assert "Edited:" in result
    edited = Document(str(path))
    texts = [p.text for p in edited.paragraphs if p.text.strip()]
    assert "Final status" in texts
    assert "Approved" in texts


@pytest.mark.asyncio
async def test_append_xlsx_appends_rows_to_sheet(tmp_path):
    from openpyxl import Workbook, load_workbook

    from miqi.documents.xlsx_tool import AppendXlsxTool

    files_dir = tmp_path / "files"
    files_dir.mkdir()
    path = files_dir / "data.xlsx"
    wb = Workbook()
    wb.active.title = "Data"
    wb.active.append(["A", "B"])
    wb.save(str(path))

    tool = AppendXlsxTool(workspace=files_dir, allowed_dir=files_dir)
    result = await tool.execute(
        filename="data.xlsx",
        sheet_name="Data",
        rows=[[1, 2], [3, 4]],
    )

    assert "Appended:" in result
    edited = load_workbook(str(path))
    assert edited["Data"]["A2"].value == 1
    assert edited["Data"]["B3"].value == 4


@pytest.mark.parametrize(
    "tool_cls, kwargs, expected_name",
    [
        pytest.param("docx", {"filename": "../escape", "content": "x"}, "escape.docx"),
        pytest.param("xlsx", {"filename": "../escape", "sheets": {}}, "escape.xlsx"),
        pytest.param("pptx", {"filename": "../escape", "slides": []}, "escape.pptx"),
    ],
)
@pytest.mark.asyncio
async def test_create_office_tools_reject_path_traversal(
    tmp_path, tool_cls, kwargs, expected_name,
):
    if tool_cls == "docx":
        from miqi.documents.docx_tool import CreateDocxTool as Tool
    elif tool_cls == "xlsx":
        from miqi.documents.xlsx_tool import CreateXlsxTool as Tool
    else:
        from miqi.documents.pptx_tool import CreatePptxTool as Tool

    files_dir = tmp_path / "files"
    tool = Tool(workspace=files_dir, allowed_dir=files_dir)

    result = await tool.execute(**kwargs)

    assert "Permission denied" in result
    assert not (tmp_path / expected_name).exists()
